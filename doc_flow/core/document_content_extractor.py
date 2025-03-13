import io
import logging
import fitz
import pptx
import openpyxl
from PIL import Image

from services.ocr_service import ocr_service

logger = logging.getLogger(__name__)

class FileProcessor:
    @staticmethod
    async def extractTextFromPdf(file_bytes: bytes, scan_images: bool = False) -> str:
        """
        Extract text from a PDF file more efficiently, and OCR images in parallel if scan_images is True.
        """
        text_lines = []
        try:
            with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                for page_num, page in enumerate(doc, start=1):
                    page_text = page.get_text("text").strip()
                    if page_text:
                        text_lines.append(f"--- Page {page_num} Text ---")
                        text_lines.append(page_text)

                    if scan_images:
                        image_list = page.get_images(full=True)
                        if not image_list:
                            continue

                        image_bytes_list = []
                        for img in image_list:
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            image_bytes_list.append(base_image["image"])

                        ocr_texts = await ocr_service.processImagesWithOcr(image_bytes_list)
                        for idx, ocr_text in enumerate(ocr_texts, start=1):
                            if ocr_text:
                                text_lines.append(f"--- Page {page_num} Image {idx} OCR ---")
                                text_lines.append(ocr_text)

        except Exception as e:
            logger.exception("PDF Extraction Error")

        return "\n".join(text_lines).replace("\t", " ")

    @staticmethod
    async def extractTextFromPptx(file_bytes: bytes, scan_images: bool = False) -> str:
        """
        Extract text from a PPTX file, and OCR images if scan_images is True.
        """
        text_lines = []
        try:
            presentation = pptx.Presentation(io.BytesIO(file_bytes))
            for slide_num, slide in enumerate(presentation.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape_text = shape.text.strip()
                        if shape_text:
                            slide_text.append(shape_text)

                    if scan_images and hasattr(shape, "image"):
                        image_bytes = shape.image.blob
                        pil_image = Image.open(io.BytesIO(image_bytes))
                        ocr_text = ocr_service.performOcr(pil_image)
                        if ocr_text:
                            slide_text.append(f"OCR Image Text: {ocr_text}")

                if slide_text:
                    text_lines.append(f"--- Slide {slide_num} ---")
                    text_lines.append("\n".join(slide_text))

        except Exception as e:
            logger.exception("PPTX Extraction Error")

        return "\n".join(text_lines).replace("\t", " ")

    @staticmethod
    async def extractTextFromXlsx(file_bytes: bytes) -> str:
        """
        Extract text from an XLSX file.
        """
        text_lines = []
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                sheet_text = []
                for row in worksheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell).strip() if cell is not None else "" for cell in row])
                    if row_text.strip():
                        sheet_text.append(row_text)
                
                if sheet_text:
                    text_lines.append(f"--- Sheet: {sheet_name} ---")
                    text_lines.append("\n".join(sheet_text))

        except Exception as e:
            logger.exception("XLSX Extraction Error")
            text_lines.append(f"XLSX Extraction Error: {str(e)}")

        return "\n".join(text_lines).replace("\t", " ")
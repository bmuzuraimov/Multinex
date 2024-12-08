import io
import logging
import os
import asyncio
from typing import List, Dict

import fitz  # PyMuPDF
import pptx
import openpyxl
import pytesseract
from PIL import Image

from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse


# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
)
logger = logging.getLogger(__name__)

# Environment-based configuration for CORS
ALLOWED_ORIGINS = os.environ.get("ALLOWED_ORIGINS", "https://typit.app,https://typeit.up.railway.app")
allowed_origins = [origin.strip() for origin in ALLOWED_ORIGINS.split(",")]

app = FastAPI(
    title="File Text Extraction API",
    description="A service to extract text from PDF, PPTX, and XLSX files, including OCR on images.",
    version="1.0.0",
    docs_url=None,
    redoc_url=None
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_credentials=True,
    allow_methods=["POST", "GET"],
    allow_headers=["*"],
)


def perform_ocr(image: Image.Image) -> str:
    """
    Perform OCR on a PIL image using Tesseract.

    :param image: PIL Image instance.
    :return: Extracted text from the image.
    """
    try:
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception as e:
        logger.error(f"OCR failed: {e}")
        return f"OCR Error: {str(e)}"


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """
    Extract text from a PDF file. Includes OCR on embedded images.

    :param file_bytes: File content in bytes.
    :return: Extracted text from the PDF.
    """
    text = ""
    try:
        with fitz.open(stream=file_bytes, filetype="pdf") as doc:
            for page_num, page in enumerate(doc, start=1):
                # Extract textual content
                page_text = page.get_text().strip()
                text += f"\n--- Page {page_num} Text ---\n{page_text}\n"

                # Extract images for OCR
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list, start=1):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = perform_ocr(image).strip()
                    if ocr_text:
                        text += f"\n--- Page {page_num} Image {img_index} OCR ---\n{ocr_text}\n"

    except Exception as e:
        logger.exception("PDF Extraction Error")
        text += f"\nPDF Extraction Error: {str(e)}\n"

    return text.replace("\t", " ").replace("\n\n", "\n")


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extract text from a PPTX file. Includes OCR on images found in slides.

    :param file_bytes: File content in bytes.
    :return: Extracted text from the PPTX file.
    """
    text = ""
    try:
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                # Extract text from shapes
                if hasattr(shape, "text"):
                    shape_text = shape.text.strip()
                    if shape_text:
                        slide_text.append(shape_text)

                # If shape is a picture (shape_type == MSO_SHAPE_TYPE.PICTURE = 13)
                # you can also confirm type by checking if shape has 'image' attribute
                if hasattr(shape, "image"):
                    image_bytes = shape.image.blob
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    ocr_text = perform_ocr(pil_image)
                    if ocr_text:
                        slide_text.append(f"OCR Image Text: {ocr_text}")

            slide_content = "\n".join(slide_text)
            text += f"\n--- Slide {slide_num} ---\n{slide_content}\n"

    except Exception as e:
        logger.exception("PPTX Extraction Error")
        text += f"\nPPTX Extraction Error: {str(e)}\n"

    return text.replace("\t", " ").replace("\n\n", "\n")


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """
    Extract text from an XLSX file.

    :param file_bytes: File content in bytes.
    :return: Extracted text from the XLSX file.
    """
    text = ""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = []
            for row in ws.iter_rows(values_only=True):
                # Convert each cell to string, strip and handle None
                row_text = "\t".join([str(cell).strip() if cell is not None else "" for cell in row])
                sheet_text.append(row_text)
            sheet_content = "\n".join(sheet_text)
            text += f"\n--- Sheet: {sheet_name} ---\n{sheet_content}\n"

    except Exception as e:
        logger.exception("XLSX Extraction Error")
        text += f"\nXLSX Extraction Error: {str(e)}\n"

    return text.replace("\t", " ").replace("\n\n", "\n")


async def process_file(filename: str, content: bytes) -> Dict[str, str]:
    """
    Asynchronously process a file based on its extension, extracting text and OCR from images if applicable.

    :param filename: Name of the file.
    :param content: Byte content of the file.
    :return: Dictionary with filename and extracted text.
    """
    extension = filename.split(".")[-1].lower()
    text = ""

    # Use asyncio.to_thread to move blocking operations off the main event loop.
    if extension == "pdf":
        text = await asyncio.to_thread(extract_text_from_pdf, content)
    elif extension in ["ppt", "pptx"]:
        text = await asyncio.to_thread(extract_text_from_pptx, content)
    elif extension in ["xls", "xlsx"]:
        text = await asyncio.to_thread(extract_text_from_xlsx, content)
    else:
        logger.warning(f"Unsupported file type requested: {extension}")
        text = "Unsupported file type."

    return {"filename": filename, "text": text}


@app.post("/extract-text")
async def extract_text(files: List[UploadFile] = File(...)) -> JSONResponse:
    """
    Endpoint to extract text from uploaded files. Supports PDF, PPTX, and XLSX files.
    Includes OCR for embedded images.

    :param files: List of files uploaded via form-data.
    :return: JSON response with extracted text.
    """
    if not files:
        raise HTTPException(status_code=400, detail="No files uploaded.")

    # Process only the first file for now, or extend logic to handle multiple files.
    first_file = files[0]
    filename = first_file.filename
    logger.info(f"Processing file: {filename}")

    try:
        content = await first_file.read()
    except Exception as e:
        logger.exception("File reading error")
        raise HTTPException(status_code=500, detail=f"Error reading file: {str(e)}")

    results = await process_file(filename, content)
    return JSONResponse(content=results)


@app.get("/")
def read_root() -> Dict[str, str]:
    """
    Root endpoint providing a welcome message and instructions.

    :return: Dictionary with a welcome message.
    """
    return {
        "message": "Welcome to the File Text Extraction API. Use the /extract-text endpoint to upload files."
    }